"""
PowerPoint export service for clinical debrief presentations.

Renders all 13+ slide types into professional PPTX format with
consistent styling, speaker notes, and clinical-friendly design.
"""
from typing import List, Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import io


# =============================================================================
# Color Palette
# =============================================================================
NAVY = RGBColor(30, 58, 95)         # #1E3A5F  — titles, headers
ACCENT_BLUE = RGBColor(37, 99, 235) # #2563EB  — links, ratings, timeline
GREEN = RGBColor(5, 150, 105)       # #059669  — strengths, achievements
AMBER = RGBColor(217, 119, 6)       # #D97706  — growth areas, warnings
RED = RGBColor(220, 38, 38)         # #DC2626  — critical badges
PURPLE = RGBColor(124, 58, 237)     # #7C3AED  — learning points, discussion
BODY_TEXT = RGBColor(31, 41, 55)    # #1F2937  — body text
SECONDARY = RGBColor(107, 114, 128) # #6B7280  — labels, captions
WHITE = RGBColor(255, 255, 255)     # #FFFFFF  — backgrounds
CARD_FILL = RGBColor(249, 250, 251) # #F9FAFB  — card backgrounds
CARD_BORDER = RGBColor(229, 231, 235)  # #E5E7EB — light borders
LIGHT_GRAY = RGBColor(243, 244, 246)   # #F3F4F6 — video block bg

# Severity-specific colors
CRITICAL_BG = RGBColor(254, 242, 242)
CRITICAL_TEXT = RGBColor(185, 28, 28)
CRITICAL_BORDER = RGBColor(254, 202, 202)
WARNING_BG = RGBColor(255, 251, 235)
WARNING_TEXT = RGBColor(180, 83, 9)
WARNING_BORDER = RGBColor(253, 230, 138)

# Rating label map
RATING_LABELS = {
    1: "Needs Improvement",
    2: "Below Average",
    3: "Average",
    4: "Above Average",
    5: "Excellent",
}

# Slide margins
MARGIN = 0.8
CONTENT_LEFT = MARGIN
CONTENT_RIGHT = 10 - MARGIN
CONTENT_WIDTH = CONTENT_RIGHT - CONTENT_LEFT


# =============================================================================
# Main Export Function
# =============================================================================

def export_to_powerpoint(session_id: str, slides: List[Dict[str, Any]]) -> bytes:
    """
    Export presentation slides to PowerPoint format.

    Args:
        session_id: The session identifier
        slides: List of slide data dictionaries

    Returns:
        Bytes of the PowerPoint file
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    total_slides = len(slides)

    # Dispatch table for slide types
    handlers = {
        'title': _add_title_slide,
        'scenario': _add_scenario_slide,
        'timeline': _add_timeline_slide,
        'moment_evidence': _add_moment_evidence_slide,
        'moment_feedback': _add_moment_feedback_slide,
        'moment': _add_moment_legacy_slide,  # backward compat
        'clinical_skill': _add_clinical_skill_slide,
        'abcde': _add_abcde_slide,
        'nts_overview': _add_nts_overview_slide,
        'nts_skill': _add_nts_skill_slide,
        'crm': _add_crm_slide,
        'pause_reflect': _add_pause_reflect_slide,
        'learning_outcomes': _add_learning_outcomes_slide,
        'action_plan': _add_action_plan_slide,
        'discussion': _add_discussion_slide,
        'summary': _add_summary_slide,
    }

    for slide_data in slides:
        slide_type = slide_data.get('type', 'moment')
        handler = handlers.get(slide_type)
        if handler:
            handler(prs, slide_data, session_id, total_slides)

    # Save to bytes
    pptx_bytes = io.BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)
    return pptx_bytes.getvalue()


# =============================================================================
# Shared Helpers
# =============================================================================

def _safe_text(value, fallback=""):
    """Return value as string, never None/null."""
    if value is None:
        return fallback
    return str(value)


def _format_time(seconds):
    """Format seconds as M:SS."""
    if seconds is None:
        return "0:00"
    seconds = float(seconds)
    mins = int(seconds // 60)
    secs = int(seconds % 60)
    return f"{mins}:{secs:02d}"


def _format_student_id(student_id: str) -> str:
    """Format student_id into readable name: 'cam1_person_1' -> 'Cam1 Person 1'."""
    if not student_id:
        return "Student"
    return student_id.replace("_", " ").title()


def _new_slide(prs):
    """Create a blank slide."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_speaker_notes(slide, notes_text: str):
    """Add speaker notes to a slide (visible in Presenter View)."""
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    notes_frame = notes_slide.notes_text_frame
    notes_frame.text = notes_text


def add_slide_footer(slide, session_id: str, slide_num: int, total: int):
    """Add consistent footer on every slide."""
    # Left: session identifier
    left_box = slide.shapes.add_textbox(
        Inches(MARGIN), Inches(7.1), Inches(4), Inches(0.3)
    )
    lf = left_box.text_frame
    lf.text = f"Clinical Simulation Debrief — {session_id}"
    lf.paragraphs[0].font.size = Pt(9)
    lf.paragraphs[0].font.color.rgb = SECONDARY

    # Right: slide number
    right_box = slide.shapes.add_textbox(
        Inches(6), Inches(7.1), Inches(3.2), Inches(0.3)
    )
    rf = right_box.text_frame
    rf.text = f"Slide {slide_num} of {total}"
    rf.paragraphs[0].font.size = Pt(9)
    rf.paragraphs[0].font.color.rgb = SECONDARY
    rf.paragraphs[0].alignment = PP_ALIGN.RIGHT


def add_section_divider(slide, y: float):
    """Add a thin blue accent line below title."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(MARGIN), Inches(y), Inches(CONTENT_WIDTH), Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()


def add_feedback_box(slide, text: str, label: str, color, left, top, width, height):
    """Add a feedback box with left accent bar and label."""
    # Left accent bar (workaround for per-side border)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, Inches(0.06), height
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = color
    accent.line.fill.background()

    # Main card
    card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left + Inches(0.06), top, width - Inches(0.06), height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_FILL
    card.line.fill.background()

    tf = card.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    tf.word_wrap = True

    # Label
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = color

    # Content
    if text:
        p2 = tf.add_paragraph()
        p2.text = _safe_text(text)
        p2.font.size = Pt(14)
        p2.font.color.rgb = BODY_TEXT
        p2.space_before = Pt(6)


def add_rating_dots(slide, rating: int, left: float, top: float):
    """Render 5 filled/empty circle dots for a 1-5 rating."""
    rating = max(0, min(5, rating or 0))
    for i in range(5):
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left + i * 0.25), Inches(top),
            Inches(0.2), Inches(0.2)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT_BLUE if i < rating else CARD_BORDER
        dot.line.fill.background()


def add_severity_badge(slide, severity: str, left: float, top: float):
    """Add a colored severity badge."""
    if not severity:
        return
    is_critical = severity.lower() == 'critical'
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(1.2), Inches(0.35)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = CRITICAL_BG if is_critical else WARNING_BG
    badge.line.color.rgb = CRITICAL_BORDER if is_critical else WARNING_BORDER
    badge.line.width = Pt(1)

    tf = badge.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_top = Inches(0.02)
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = severity.upper()
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = CRITICAL_TEXT if is_critical else WARNING_TEXT
    p.alignment = PP_ALIGN.CENTER


def add_checklist_row(slide, label: str, checked, left: float, top: float):
    """Add a CRM checklist row with Unicode check/cross/dash."""
    if checked is True:
        icon = "\u2713"
        icon_color = GREEN
    elif checked is False:
        icon = "\u2717"
        icon_color = RED
    else:
        icon = "\u2014"
        icon_color = SECONDARY

    # Icon
    icon_box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(0.4), Inches(0.35)
    )
    ip = icon_box.text_frame.paragraphs[0]
    ip.text = icon
    ip.font.size = Pt(16)
    ip.font.bold = True
    ip.font.color.rgb = icon_color
    ip.alignment = PP_ALIGN.CENTER

    # Label
    label_box = slide.shapes.add_textbox(
        Inches(left + 0.45), Inches(top), Inches(3.5), Inches(0.35)
    )
    lp = label_box.text_frame.paragraphs[0]
    lp.text = label
    lp.font.size = Pt(14)
    lp.font.color.rgb = BODY_TEXT


def _add_title_text(slide, title: str, top: float = 0.8, centered: bool = False):
    """Add a standard slide title."""
    title_box = slide.shapes.add_textbox(
        Inches(MARGIN), Inches(top), Inches(CONTENT_WIDTH), Inches(0.8)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = _safe_text(title, "Untitled")
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = NAVY
    if centered:
        p.alignment = PP_ALIGN.CENTER


def _add_video_reference_block(slide, video_clip: dict, left: float, top: float,
                                width: float = 4.2, height: float = 3.5):
    """Add a video reference placeholder block."""
    block = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    block.fill.solid()
    block.fill.fore_color.rgb = LIGHT_GRAY
    block.line.fill.background()

    tf = block.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)

    start = _format_time(video_clip.get('start', 0))
    end = _format_time(video_clip.get('end', 0))

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(40)
    run = p.add_run()
    run.text = f"Video: {start} - {end}"
    run.font.size = Pt(18)
    run.font.color.rgb = SECONDARY
    run.font.bold = True

    cameras = video_clip.get('cameras', ['cam1', 'cam2', 'cam3', 'monitor'])
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)
    run2 = p2.add_run()
    run2.text = f"Refer to session recording\nCameras: {', '.join(cameras)}"
    run2.font.size = Pt(11)
    run2.font.color.rgb = SECONDARY


def _add_metric_box(slide, value, label: str, bg_color, left, top, width=2.0, height=1.1):
    """Add a metric box with large number and label."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.15)

    p = tf.paragraphs[0]
    p.text = str(value)
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BODY_TEXT

    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(10)
    p2.font.bold = True
    p2.font.color.rgb = SECONDARY
    p2.space_before = Pt(6)


def _add_numbered_badge(slide, number: int, left: float, top: float, color=None):
    """Add a numbered circle badge."""
    badge_color = color or ACCENT_BLUE
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(0.4), Inches(0.4)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = badge_color
    circle.line.fill.background()

    # Number text (on top of circle)
    num_box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(0.4), Inches(0.4)
    )
    tf = num_box.text_frame
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


# =============================================================================
# Slide Type Handlers
# =============================================================================

def _add_title_slide(prs, slide_data, session_id, total_slides):
    """Slide 1: Title — session overview."""
    slide = _new_slide(prs)
    slide_num = slide_data.get('slide_number', 1)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.0), Inches(8), Inches(1.2)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = _safe_text(slide_data.get('title'), 'Clinical Simulation Debrief')
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    # Subtitle (student name)
    subtitle = slide_data.get('subtitle', '')
    content = slide_data.get('content', {})
    metadata = slide_data.get('metadata', {})

    display_name = subtitle or metadata.get('student', '')
    if not display_name or display_name == 'null':
        student_id = content.get('student_id', '') or metadata.get('student_id', '')
        display_name = _format_student_id(student_id) if student_id else 'Student'

    sub_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(3.4), Inches(7), Inches(0.8)
    )
    sf = sub_box.text_frame
    sp = sf.paragraphs[0]
    sp.text = display_name
    sp.font.size = Pt(28)
    sp.font.color.rgb = ACCENT_BLUE
    sp.font.bold = True
    sp.alignment = PP_ALIGN.CENTER

    # Role + Scenario type tag
    role = content.get('role', '') or metadata.get('role', '')
    scenario_type = content.get('scenario_type', '') or metadata.get('scenario_type', '')
    tag_parts = [p for p in [role, scenario_type] if p]
    if tag_parts:
        tag_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(4.2), Inches(7), Inches(0.4)
        )
        tp = tag_box.text_frame.paragraphs[0]
        tp.text = " | ".join(tag_parts)
        tp.font.size = Pt(16)
        tp.font.color.rgb = SECONDARY
        tp.alignment = PP_ALIGN.CENTER

    # Divider line
    add_section_divider(slide, 5.0)

    # Metadata row: Student / Evaluator / Session
    y = 5.3
    cols = []
    student_name = metadata.get('student', display_name)
    evaluator_name = metadata.get('evaluator', '')
    if student_name:
        cols.append(("STUDENT", student_name))
    if evaluator_name:
        cols.append(("EVALUATOR", evaluator_name))
    cols.append(("SESSION", session_id))

    col_width = CONTENT_WIDTH / max(len(cols), 1)
    for i, (lbl, val) in enumerate(cols):
        x = MARGIN + i * col_width
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(col_width), Inches(0.9))
        btf = box.text_frame
        bp = btf.paragraphs[0]
        bp.text = lbl
        bp.font.size = Pt(10)
        bp.font.bold = True
        bp.font.color.rgb = SECONDARY

        bp2 = btf.add_paragraph()
        bp2.text = _safe_text(val)
        bp2.font.size = Pt(18)
        bp2.font.color.rgb = BODY_TEXT
        bp2.space_before = Pt(4)

    add_slide_footer(slide, session_id, slide_num, total_slides)
    add_speaker_notes(slide, "Welcome to the debrief session. Review objectives before proceeding.")


def _add_scenario_slide(prs, slide_data, session_id, total_slides):
    """Slide 2: Scenario Context — patient presentation + learning objectives."""
    slide = _new_slide(prs)
    slide_num = slide_data.get('slide_number', 2)

    _add_title_text(slide, slide_data.get('title', 'Scenario Context'))
    add_section_divider(slide, 1.6)

    content = slide_data.get('content', {})
    left_x = MARGIN
    right_x = 5.2

    # Patient presentation (blue box)
    patient = content.get('patient_presentation', '')
    if patient:
        add_feedback_box(
            slide, patient, "PATIENT PRESENTATION", ACCENT_BLUE,
            Inches(left_x), Inches(1.9), Inches(4.0), Inches(1.6)
        )

    # Clinical context (gray box)
    context = content.get('clinical_context', '')
    if context:
        add_feedback_box(
            slide, context, "CLINICAL CONTEXT", SECONDARY,
            Inches(left_x), Inches(3.7), Inches(4.0), Inches(1.4)
        )

    # Learning objectives
    objectives = content.get('learning_objectives', [])
    if objectives:
        obj_header = slide.shapes.add_textbox(
            Inches(right_x), Inches(1.9), Inches(4.0), Inches(0.4)
        )
        oh = obj_header.text_frame.paragraphs[0]
        oh.text = "LEARNING OBJECTIVES"
        oh.font.size = Pt(10)
        oh.font.bold = True
        oh.font.color.rgb = SECONDARY

        y_pos = 2.4
        for idx, obj in enumerate(objectives[:6], 1):
            _add_numbered_badge(slide, idx, right_x, y_pos)
            obj_box = slide.shapes.add_textbox(
                Inches(right_x + 0.5), Inches(y_pos), Inches(3.5), Inches(0.5)
            )
            op = obj_box.text_frame.paragraphs[0]
            op.text = _safe_text(obj)
            op.font.size = Pt(14)
            op.font.color.rgb = BODY_TEXT
            obj_box.text_frame.word_wrap = True
            y_pos += 0.55

    add_slide_footer(slide, session_id, slide_num, total_slides)
    add_speaker_notes(
        slide,
        "Review the patient presentation with the student. "
        "Confirm understanding of learning objectives. "
        "Set expectations for the debrief discussion."
    )


def _add_timeline_slide(prs, slide_data, session_id, total_slides):
    """Slide 3: Session Timeline Overview — moments plotted on a horizontal bar."""
    slide = _new_slide(prs)
    slide_num = slide_data.get('slide_number', 3)

    _add_title_text(slide, slide_data.get('title', 'Session Timeline'))
    add_section_divider(slide, 1.6)

    moments = slide_data.get('content', {}).get('moments', [])
    if not moments:
        no_data = slide.shapes.add_textbox(
            Inches(MARGIN), Inches(3), Inches(CONTENT_WIDTH), Inches(1)
        )
        no_data.text_frame.paragraphs[0].text = "No critical moments detected in this session."
        no_data.text_frame.paragraphs[0].font.size = Pt(18)
        no_data.text_frame.paragraphs[0].font.color.rgb = SECONDARY
        no_data.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        add_slide_footer(slide, session_id, slide_num, total_slides)
        return

    # Calculate timeline bounds
    session_end = max(m.get('end_time', 0) for m in moments)
    if session_end <= 0:
        session_end = 600  # fallback 10 min

    timeline_left = MARGIN + 0.3
    timeline_width = CONTENT_WIDTH - 0.6
    timeline_y = 3.2

    # Timeline baseline
    baseline = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(timeline_left), Inches(timeline_y),
        Inches(timeline_width), Inches(0.03)
    )
    baseline.fill.solid()
    baseline.fill.fore_color.rgb = CARD_BORDER
    baseline.line.fill.background()

    # Start/end labels
    start_lbl = slide.shapes.add_textbox(
        Inches(timeline_left), Inches(timeline_y + 0.1), Inches(0.5), Inches(0.3)
    )
    start_lbl.text_frame.paragraphs[0].text = "0:00"
    start_lbl.text_frame.paragraphs[0].font.size = Pt(9)
    start_lbl.text_frame.paragraphs[0].font.color.rgb = SECONDARY

    end_lbl = slide.shapes.add_textbox(
        Inches(timeline_left + timeline_width - 0.5), Inches(timeline_y + 0.1),
        Inches(0.5), Inches(0.3)
    )
    end_lbl.text_frame.paragraphs[0].text = _format_time(session_end)
    end_lbl.text_frame.paragraphs[0].font.size = Pt(9)
    end_lbl.text_frame.paragraphs[0].font.color.rgb = SECONDARY
    end_lbl.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    # Plot moments (max 10)
    display_moments = sorted(moments, key=lambda m: m.get('start_time', 0))[:10]
    for m in display_moments:
        start_time = m.get('start_time', 0)
        severity = m.get('severity', 'warning')
        moment_id = m.get('id', '?')
        summary = _safe_text(m.get('summary', ''))[:60]

        proportion = start_time / session_end if session_end > 0 else 0
        marker_x = timeline_left + proportion * timeline_width

        # Marker dot
        marker = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(marker_x - 0.1), Inches(timeline_y - 0.1),
            Inches(0.25), Inches(0.25)
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = RED if severity == 'critical' else AMBER
        marker.line.fill.background()

        # Label below
        label_box = slide.shapes.add_textbox(
            Inches(max(timeline_left, marker_x - 0.5)),
            Inches(timeline_y + 0.5),
            Inches(1.2), Inches(0.8)
        )
        ltf = label_box.text_frame
        ltf.word_wrap = True

        lp1 = ltf.paragraphs[0]
        lp1.text = f"#{moment_id}"
        lp1.font.size = Pt(10)
        lp1.font.bold = True
        lp1.font.color.rgb = RED if severity == 'critical' else AMBER

        if summary:
            lp2 = ltf.add_paragraph()
            lp2.text = summary
            lp2.font.size = Pt(8)
            lp2.font.color.rgb = SECONDARY

    if len(moments) > 10:
        extra = slide.shapes.add_textbox(
            Inches(MARGIN), Inches(6.2), Inches(CONTENT_WIDTH), Inches(0.4)
        )
        extra.text_frame.paragraphs[0].text = f"+{len(moments) - 10} additional moments not shown"
        extra.text_frame.paragraphs[0].font.size = Pt(11)
        extra.text_frame.paragraphs[0].font.color.rgb = SECONDARY
        extra.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    add_slide_footer(slide, session_id, slide_num, total_slides)
    add_speaker_notes(
        slide,
        "Select which moments to focus on. You can skip less critical ones."
    )


def _add_moment_evidence_slide(prs, slide_data, session_id, total_slides):
    """Moment Evidence slide — objective, multi-source evidence."""
    slide = _new_slide(prs)
    slide_num = slide_data.get('slide_number', 0)
    severity = slide_data.get('severity', '')
    video_clip = slide_data.get('video_clip', {})

    # Severity badge
    add_severity_badge(slide, severity, MARGIN, 0.8)

    # Timestamp
    if video_clip:
        time_box = slide.shapes.add_textbox(
            Inches(MARGIN + 1.4), Inches(0.85), Inches(2), Inches(0.3)
        )
        tp = time_box.text_frame.paragraphs[0]
        tp.text = f"{_format_time(video_clip.get('start', 0))} - {_format_time(video_clip.get('end', 0))}"
        tp.font.size = Pt(12)
        tp.font.color.rgb = SECONDARY

    # Title
    _add_title_text(slide, slide_data.get('title', 'Critical Moment'), top=1.3)

    # Left: Video reference block
    if video_clip:
        _add_video_reference_block(slide, video_clip, MARGIN, 2.3, width=4.2, height=3.2)

    # Right: Evidence items
    observations = slide_data.get('key_observations', [])
    right_x = 5.3
    y_pos = 2.3
    for obs in observations[:4]:
        obs_text = _safe_text(obs)
        # Determine source tag from prefix like "[video_cam1]"
        source_tag = ""
        if obs_text.startswith("["):
            bracket_end = obs_text.find("]")
            if bracket_end > 0:
                source_tag = obs_text[1:bracket_end]
                obs_text = obs_text[bracket_end + 1:].strip()

        obs_box = slide.shapes.add_textbox(
            Inches(right_x), Inches(y_pos), Inches(4.0), Inches(0.7)
        )
        otf = obs_box.text_frame
        otf.word_wrap = True
        op = otf.paragraphs[0]

        if source_tag:
            tag_run = op.add_run()
            tag_run.text = f"[{source_tag}] "
            tag_run.font.size = Pt(11)
            tag_run.font.bold = True
            if 'video' in source_tag.lower():
                tag_run.font.color.rgb = ACCENT_BLUE
            elif 'transcript' in source_tag.lower():
                tag_run.font.color.rgb = PURPLE
            elif 'monitor' in source_tag.lower():
                tag_run.font.color.rgb = RED
            else:
                tag_run.font.color.rgb = SECONDARY

        text_run = op.add_run()
        text_run.text = obs_text
        text_run.font.size = Pt(13)
        text_run.font.color.rgb = BODY_TEXT

        y_pos += 0.8

    # Transcript snippets (bottom strip)
    transcript = slide_data.get('transcript_snippets', [])
    if transcript:
        ts_box = slide.shapes.add_textbox(
            Inches(MARGIN), Inches(5.7), Inches(CONTENT_WIDTH), Inches(0.9)
        )
        ts_tf = ts_box.text_frame
        ts_tf.word_wrap = True
        for seg in transcript[:2]:
            speaker = _safe_text(seg.get('speaker', ''))
            text = _safe_text(seg.get('text', ''))
            p = ts_tf.paragraphs[0] if not ts_tf.paragraphs[0].text else ts_tf.add_paragraph()
            if speaker:
                sr = p.add_run()
                sr.text = f"{speaker}: "
                sr.font.bold = True
                sr.font.size = Pt(14)
                sr.font.color.rgb = ACCENT_BLUE
            tr = p.add_run()
            tr.text = text
            tr.font.size = Pt(14)
            tr.font.italic = True
            tr.font.color.rgb = BODY_TEXT

    add_slide_footer(slide, session_id, slide_num, total_slides)
    add_speaker_notes(
        slide,
        "Ask: 'What do you think happened here?' before advancing to the feedback slide."
    )


def _add_moment_feedback_slide(prs, slide_data, session_id, total_slides):
    """Moment Feedback slide — evaluator feedback boxes."""
    slide = _new_slide(prs)
    slide_num = slide_data.get('slide_number', 0)
    severity = slide_data.get('severity', '')
    video_clip = slide_data.get('video_clip', {})

    # Severity badge + timestamp header
    add_severity_badge(slide, severity, MARGIN, 0.8)
    if video_clip:
        time_box = slide.shapes.add_textbox(
            Inches(MARGIN + 1.4), Inches(0.85), Inches(2), Inches(0.3)
        )
        tp = time_box.text_frame.paragraphs[0]
        tp.text = f"{_format_time(video_clip.get('start', 0))} - {_format_time(video_clip.get('end', 0))}"
        tp.font.size = Pt(12)
        tp.font.color.rgb = SECONDARY

    _add_title_text(slide, slide_data.get('title', 'Feedback'), top=1.3)

    # Feedback boxes — full width, stacked
    content = slide_data.get('presenter_notes', slide_data.get('content', {}))
    box_left = MARGIN
    box_width = CONTENT_WIDTH
    y_pos = 2.3

    # Clinical Context
    clinical = content.get('clinical_significance', '')
    if clinical:
        add_feedback_box(slide, clinical, "CLINICAL CONTEXT", ACCENT_BLUE,
                         Inches(box_left), Inches(y_pos), Inches(box_width), Inches(1.1))
        y_pos += 1.3

    # Key Strength
    positive = content.get('positive_feedback', '')
    if positive:
        add_feedback_box(slide, positive, "KEY STRENGTH", GREEN,
                         Inches(box_left), Inches(y_pos), Inches(box_width), Inches(1.1))
        y_pos += 1.3

    # Area for Growth
    negative = content.get('areas_for_improvement', '')
    if negative:
        add_feedback_box(slide, negative, "AREA FOR GROWTH", AMBER,
                         Inches(box_left), Inches(y_pos), Inches(box_width), Inches(1.1))
        y_pos += 1.3

    # Learning Point
    learning = content.get('learning_points', '')
    if learning:
        add_feedback_box(slide, f'"{learning}"', "KEY TAKEAWAY", PURPLE,
                         Inches(box_left), Inches(y_pos), Inches(box_width), Inches(1.1))

    add_slide_footer(slide, session_id, slide_num, total_slides)
    add_speaker_notes(
        slide,
        "Discuss each feedback point. Encourage the student to reflect on their actions."
    )


def _add_moment_legacy_slide(prs, slide_data, session_id, total_slides):
    """Backward-compat handler for the old combined 'moment' type."""
    # Split into evidence + feedback if presenter_notes exist
    has_feedback = False
    notes = slide_data.get('presenter_notes', slide_data.get('content', {}))
    if notes:
        has_feedback = any(notes.get(k) for k in
                          ['positive_feedback', 'areas_for_improvement',
                           'learning_points', 'clinical_significance'])

    _add_moment_evidence_slide(prs, slide_data, session_id, total_slides)
    if has_feedback:
        _add_moment_feedback_slide(prs, slide_data, session_id, total_slides)


def _add_clinical_skill_slide(prs, slide_data, session_id, total_slides):
    """Clinical Skill slide — per-skill deep-dive with rating and observations."""
    slide = _new_slide(prs)
    slide_num = slide_data.get('slide_number', 0)
    content = slide_data.get('content', {})

    # Title + protocol reference
    _add_title_text(slide, slide_data.get('title', 'Clinical Skill'))

    subtitle = slide_data.get('subtitle', '')
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(MARGIN), Inches(1.5), Inches(4), Inches(0.3)
        )
        sp = sub_box.text_frame.paragraphs[0]
        sp.text = _safe_text(subtitle)
        sp.font.size = Pt(12)
        sp.font.color.rgb = SECONDARY

    add_section_divider(slide, 1.85)

    # Left column: rating + video
    rating = content.get('rating', content.get('performance_rating', 0))
    add_rating_dots(slide, rating, MARGIN, 2.2)

    rating_label = RATING_LABELS.get(rating, "Not Rated")
    rl_box = slide.shapes.add_textbox(
        Inches(MARGIN + 1.4), Inches(2.15), Inches(2), Inches(0.3)
    )
    rl_box.text_frame.paragraphs[0].text = rating_label
    rl_box.text_frame.paragraphs[0].font.size = Pt(14)
    rl_box.text_frame.paragraphs[0].font.bold = True
    rl_box.text_frame.paragraphs[0].font.color.rgb = (
        GREEN if rating >= 4 else AMBER if rating >= 3 else RED
    )

    # Video reference
    video_clip = slide_data.get('video_clip')
    if video_clip:
        _add_video_reference_block(slide, video_clip, MARGIN, 2.8, width=4.0, height=2.8)

    # Right column: summary + observations
    right_x = 5.2
    y_pos = 2.2

    summary = content.get('summary', '')
    if summary:
        add_feedback_box(slide, summary, "SUMMARY", ACCENT_BLUE,
                         Inches(right_x), Inches(y_pos), Inches(4.0), Inches(1.0))
        y_pos += 1.2

    observations = content.get('observations', [])
    if observations:
        obs_header = slide.shapes.add_textbox(
            Inches(right_x), Inches(y_pos), Inches(4.0), Inches(0.3)
        )
        obs_header.text_frame.paragraphs[0].text = "OBSERVATIONS"
        obs_header.text_frame.paragraphs[0].font.size = Pt(10)
        obs_header.text_frame.paragraphs[0].font.bold = True
        obs_header.text_frame.paragraphs[0].font.color.rgb = SECONDARY
        y_pos += 0.35

        for obs in observations[:4]:
            desc = _safe_text(obs.get('description', ''))
            timestamp = _safe_text(obs.get('timestamp', ''))
            source = _safe_text(obs.get('evidence_source', ''))

            obs_box = slide.shapes.add_textbox(
                Inches(right_x), Inches(y_pos), Inches(4.0), Inches(0.65)
            )
            otf = obs_box.text_frame
            otf.word_wrap = True

            op = otf.paragraphs[0]
            op.text = desc
            op.font.size = Pt(13)
            op.font.color.rgb = BODY_TEXT

            meta_p = otf.add_paragraph()
            meta_run = meta_p.add_run()
            meta_run.text = f"{timestamp}  "
            meta_run.font.size = Pt(10)
            meta_run.font.color.rgb = SECONDARY

            if source:
                src_run = meta_p.add_run()
                src_run.text = f"[{source}]"
                src_run.font.size = Pt(10)
                src_run.font.bold = True
                src_run.font.color.rgb = ACCENT_BLUE

            y_pos += 0.75

    add_slide_footer(slide, session_id, slide_num, total_slides)
    add_speaker_notes(
        slide,
        f"~2 min | Performance rating: {rating_label}. {summary}\n"
        "Discussion: What were you thinking at this moment? What would you do differently?"
    )


def _add_abcde_slide(prs, slide_data, session_id, total_slides):
    """ABCDE Assessment slide — structured primary survey table."""
    slide = _new_slide(prs)
    slide_num = slide_data.get('slide_number', 0)
    content = slide_data.get('content', {})

    _add_title_text(slide, slide_data.get('title', 'ABCDE Assessment'))

    # Sequence badge
    seq_followed = content.get('overall_sequence_followed')
    if seq_followed is not None:
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.5), Inches(0.85), Inches(2.5), Inches(0.4)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(220, 252, 231) if seq_followed else RGBColor(254, 243, 199)
        badge.line.fill.background()

        bp = badge.text_frame.paragraphs[0]
        bp.text = "\u2713 Sequence Followed" if seq_followed else "\u26A0 Sequence Incomplete"
        bp.font.size = Pt(11)
        bp.font.bold = True
        bp.font.color.rgb = GREEN if seq_followed else AMBER
        bp.alignment = PP_ALIGN.CENTER

    add_section_divider(slide, 1.6)

    # ABCDE table
    components = [
        ('A — Airway', content.get('airway', {})),
        ('B — Breathing', content.get('breathing', {})),
        ('C — Circulation', content.get('circulation', {})),
        ('D — Disability', content.get('disability', {})),
        ('E — Exposure', content.get('exposure', {})),
    ]

    # Create table: header + 5 rows, 4 columns
    table_shape = slide.shapes.add_table(
        6, 4,
        Inches(MARGIN), Inches(1.9),
        Inches(CONTENT_WIDTH), Inches(3.5)
    )
    table = table_shape.table

    # Header row
    headers = ['Component', 'Rating', 'Status', 'Actions Taken']
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY

    # Data rows
    for row_idx, (label, comp) in enumerate(components, 1):
        # Component name
        cell0 = table.cell(row_idx, 0)
        cell0.text = label
        cell0.text_frame.paragraphs[0].font.size = Pt(13)
        cell0.text_frame.paragraphs[0].font.bold = True
        cell0.text_frame.paragraphs[0].font.color.rgb = BODY_TEXT

        # Rating
        rating = comp.get('rating', 0) or 0
        cell1 = table.cell(row_idx, 1)
        dots = "\u25CF " * rating + "\u25CB " * (5 - rating)
        cell1.text = dots.strip()
        cell1.text_frame.paragraphs[0].font.size = Pt(12)
        cell1.text_frame.paragraphs[0].font.color.rgb = ACCENT_BLUE

        # Status
        cell2 = table.cell(row_idx, 2)
        cell2.text = _safe_text(comp.get('status', '—'))
        cell2.text_frame.paragraphs[0].font.size = Pt(12)
        cell2.text_frame.paragraphs[0].font.color.rgb = BODY_TEXT

        # Actions
        actions = comp.get('actions_taken', [])
        cell3 = table.cell(row_idx, 3)
        cell3.text = ", ".join(actions) if actions else "—"
        cell3.text_frame.paragraphs[0].font.size = Pt(11)
        cell3.text_frame.paragraphs[0].font.color.rgb = BODY_TEXT
        cell3.text_frame.word_wrap = True

        # Alternating row color
        bg = CARD_FILL if row_idx % 2 == 1 else WHITE
        for col in range(4):
            table.cell(row_idx, col).fill.solid()
            table.cell(row_idx, col).fill.fore_color.rgb = bg

    # Time to complete footer
    ttc = content.get('time_to_complete')
    if ttc:
        ttc_box = slide.shapes.add_textbox(
            Inches(MARGIN), Inches(5.6), Inches(CONTENT_WIDTH), Inches(0.4)
        )
        ttc_p = ttc_box.text_frame.paragraphs[0]
        ttc_p.text = f"Time to complete: {_format_time(ttc)}"
        ttc_p.font.size = Pt(12)
        ttc_p.font.color.rgb = SECONDARY

    add_slide_footer(slide, session_id, slide_num, total_slides)
    add_speaker_notes(
        slide,
        "~3 min | Walk me through your assessment process. "
        "What findings guided your priorities?"
    )


def _add_nts_overview_slide(prs, slide_data, session_id, total_slides):
    """Non-Technical Skills Overview — 5-dimension dashboard."""
    slide = _new_slide(prs)
    slide_num = slide_data.get('slide_number', 0)
    content = slide_data.get('content', {})

    _add_title_text(slide, slide_data.get('title', 'Non-Technical Skills'))
    add_section_divider(slide, 1.6)

    # Overall rating badge
    overall = content.get('overall_nts_rating')
    if overall:
        ov_box = slide.shapes.add_textbox(
            Inches(7), Inches(0.85), Inches(2.2), Inches(0.4)
        )
        ov_p = ov_box.text_frame.paragraphs[0]
        ov_p.text = f"Overall: {RATING_LABELS.get(overall, str(overall))}"
        ov_p.font.size = Pt(14)
        ov_p.font.bold = True
        ov_p.font.color.rgb = GREEN if overall >= 4 else AMBER if overall >= 3 else RED
        ov_p.alignment = PP_ALIGN.RIGHT

    # 5 columns
    skills = [
        ('Communication', content.get('communication', {})),
        ('Teamwork', content.get('teamwork', {})),
        ('Leadership', content.get('leadership', {})),
        ('Situational\nAwareness', content.get('situational_awareness', {})),
        ('Decision\nMaking', content.get('decision_making', {})),
    ]

    col_width = CONTENT_WIDTH / 5
    for i, (name, data) in enumerate(skills):
        x = MARGIN + i * col_width
        y_start = 2.2

        # Skill name
        name_box = slide.shapes.add_textbox(
            Inches(x), Inches(y_start), Inches(col_width), Inches(0.7)
        )
        np = name_box.text_frame.paragraphs[0]
        np.text = name
        np.font.size = Pt(14)
        np.font.bold = True
        np.font.color.rgb = BODY_TEXT
        np.alignment = PP_ALIGN.CENTER
        name_box.text_frame.word_wrap = True

        # Rating dots
        rating = data.get('rating', 0) or 0
        dot_x = x + (col_width - 1.2) / 2  # center the dots
        add_rating_dots(slide, rating, dot_x, y_start + 0.8)

        # Rating label
        lbl_box = slide.shapes.add_textbox(
            Inches(x), Inches(y_start + 1.1), Inches(col_width), Inches(0.3)
        )
        lp = lbl_box.text_frame.paragraphs[0]
        lp.text = RATING_LABELS.get(rating, "—")
        lp.font.size = Pt(10)
        lp.font.color.rgb = SECONDARY
        lp.alignment = PP_ALIGN.CENTER

    add_slide_footer(slide, session_id, slide_num, total_slides)
    add_speaker_notes(
        slide,
        "Overview slide — detailed discussion on individual skill slides that follow."
    )


def _add_nts_skill_slide(prs, slide_data, session_id, total_slides):
    """Individual NTS Skill slide — strengths + areas for improvement."""
    slide = _new_slide(prs)
    slide_num = slide_data.get('slide_number', 0)
    content = slide_data.get('content', {})

    _add_title_text(slide, slide_data.get('title', 'Skill'))

    # Rating badge
    rating = content.get('rating', 0) or 0
    add_rating_dots(slide, rating, MARGIN, 1.5)

    lbl_box = slide.shapes.add_textbox(
        Inches(MARGIN + 1.4), Inches(1.45), Inches(2.5), Inches(0.3)
    )
    lp = lbl_box.text_frame.paragraphs[0]
    lp.text = RATING_LABELS.get(rating, "—")
    lp.font.size = Pt(14)
    lp.font.bold = True
    lp.font.color.rgb = GREEN if rating >= 4 else AMBER if rating >= 3 else RED

    add_section_divider(slide, 1.9)

    # Left column: Strengths + Areas for Improvement
    left_x = MARGIN
    y_pos = 2.2

    strengths = content.get('strengths', [])
    if strengths:
        add_feedback_box(
            slide, "\n".join(f"\u2022 {s}" for s in strengths),
            "STRENGTHS", GREEN,
            Inches(left_x), Inches(y_pos), Inches(4.0), Inches(min(1.0 + 0.3 * len(strengths), 2.5))
        )
        y_pos += min(1.2 + 0.3 * len(strengths), 2.8)

    areas = content.get('areas_for_improvement', [])
    if areas:
        add_feedback_box(
            slide, "\n".join(f"\u2022 {a}" for a in areas),
            "AREAS FOR IMPROVEMENT", AMBER,
            Inches(left_x), Inches(y_pos), Inches(4.0), Inches(min(1.0 + 0.3 * len(areas), 2.0))
        )

    # Right column: Observations
    observations = content.get('observations', [])
    if observations:
        right_x = 5.2
        obs_header = slide.shapes.add_textbox(
            Inches(right_x), Inches(2.2), Inches(4.0), Inches(0.3)
        )
        obs_header.text_frame.paragraphs[0].text = "OBSERVATIONS"
        obs_header.text_frame.paragraphs[0].font.size = Pt(10)
        obs_header.text_frame.paragraphs[0].font.bold = True
        obs_header.text_frame.paragraphs[0].font.color.rgb = SECONDARY

        obs_y = 2.6
        for obs in observations[:4]:
            desc = _safe_text(obs.get('description', ''))
            ts = _safe_text(obs.get('timestamp', ''))
            obs_box = slide.shapes.add_textbox(
                Inches(right_x), Inches(obs_y), Inches(4.0), Inches(0.6)
            )
            otf = obs_box.text_frame
            otf.word_wrap = True
            otf.paragraphs[0].text = desc
            otf.paragraphs[0].font.size = Pt(13)
            otf.paragraphs[0].font.color.rgb = BODY_TEXT

            if ts:
                meta = otf.add_paragraph()
                meta.text = ts
                meta.font.size = Pt(10)
                meta.font.color.rgb = SECONDARY

            obs_y += 0.7

    skill_name = slide_data.get('title', 'this skill').lower()
    add_slide_footer(slide, session_id, slide_num, total_slides)
    add_speaker_notes(
        slide,
        f"~1.5 min | How did you approach {skill_name} during this scenario? "
        "What challenges did you face?"
    )


def _add_crm_slide(prs, slide_data, session_id, total_slides):
    """CRM Reflection Checklist slide."""
    slide = _new_slide(prs)
    slide_num = slide_data.get('slide_number', 0)
    content = slide_data.get('content', {})

    _add_title_text(slide, slide_data.get('title', 'Crisis Resource Management'))
    add_section_divider(slide, 1.6)

    # Left: CRM Principles checklist
    principles = [
        ('knew_environment', 'Knew the environment'),
        ('mobilized_resources', 'Mobilized resources'),
        ('prevented_fixation', 'Prevented fixation errors'),
        ('cross_checked', 'Cross-checked'),
        ('used_cognitive_aids', 'Used cognitive aids'),
        ('re_evaluated', 'Re-evaluated regularly'),
        ('set_priorities', 'Set priorities dynamically'),
    ]

    header_box = slide.shapes.add_textbox(
        Inches(MARGIN), Inches(1.9), Inches(4.0), Inches(0.3)
    )
    header_box.text_frame.paragraphs[0].text = "CRM PRINCIPLES"
    header_box.text_frame.paragraphs[0].font.size = Pt(10)
    header_box.text_frame.paragraphs[0].font.bold = True
    header_box.text_frame.paragraphs[0].font.color.rgb = SECONDARY

    y_pos = 2.3
    for key, label in principles:
        checked = content.get(key)
        add_checklist_row(slide, label, checked, MARGIN, y_pos)
        y_pos += 0.45

    # Right: Strengths + Areas for Development
    right_x = 5.2
    r_y = 1.9

    crm_strengths = content.get('strengths', [])
    if crm_strengths:
        add_feedback_box(
            slide, "\n".join(f"\u2022 {s}" for s in crm_strengths),
            "STRENGTHS", GREEN,
            Inches(right_x), Inches(r_y), Inches(4.0),
            Inches(min(1.0 + 0.3 * len(crm_strengths), 2.5))
        )
        r_y += min(1.2 + 0.3 * len(crm_strengths), 2.8)

    crm_areas = content.get('areas_for_development', [])
    if crm_areas:
        add_feedback_box(
            slide, "\n".join(f"\u2022 {a}" for a in crm_areas),
            "AREAS FOR DEVELOPMENT", AMBER,
            Inches(right_x), Inches(r_y), Inches(4.0),
            Inches(min(1.0 + 0.3 * len(crm_areas), 2.0))
        )

    add_slide_footer(slide, session_id, slide_num, total_slides)
    add_speaker_notes(
        slide,
        "~2 min | Focus on team dynamics and resource utilization. "
        "Discuss the 10-for-10 principle if applicable."
    )


def _add_pause_reflect_slide(prs, slide_data, session_id, total_slides):
    """Pause & Reflect slide — discussion prompt mid-deck."""
    slide = _new_slide(prs)
    slide_num = slide_data.get('slide_number', 0)
    content = slide_data.get('content', {})

    # Large centered question
    question = content.get('question', 'What would you do differently?')
    q_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.0), Inches(7), Inches(1.5)
    )
    qtf = q_box.text_frame
    qtf.word_wrap = True
    qp = qtf.paragraphs[0]
    qp.text = f'"{question}"'
    qp.font.size = Pt(32)
    qp.font.bold = True
    qp.font.color.rgb = NAVY
    qp.alignment = PP_ALIGN.CENTER

    # Options as rounded boxes
    options = content.get('options', [
        "What went well?",
        "What would you change?",
        "What will you take forward?"
    ])

    num_options = min(len(options), 3)
    total_option_width = num_options * 2.5 + (num_options - 1) * 0.3
    start_x = (10 - total_option_width) / 2

    for i, option in enumerate(options[:3]):
        x = start_x + i * 2.8
        opt_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(4.2), Inches(2.5), Inches(1.0)
        )
        opt_shape.fill.solid()
        opt_shape.fill.fore_color.rgb = CARD_FILL
        opt_shape.line.color.rgb = PURPLE
        opt_shape.line.width = Pt(2)

        otf = opt_shape.text_frame
        otf.word_wrap = True
        otf.margin_left = Inches(0.15)
        otf.margin_right = Inches(0.15)
        otf.margin_top = Inches(0.1)
        op = otf.paragraphs[0]
        op.text = _safe_text(option)
        op.font.size = Pt(14)
        op.font.color.rgb = PURPLE
        op.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(
        Inches(2), Inches(5.8), Inches(6), Inches(0.4)
    )
    sub_p = sub_box.text_frame.paragraphs[0]
    sub_p.text = "Pause for discussion before continuing"
    sub_p.font.size = Pt(14)
    sub_p.font.italic = True
    sub_p.font.color.rgb = SECONDARY
    sub_p.alignment = PP_ALIGN.CENTER

    add_slide_footer(slide, session_id, slide_num, total_slides)
    add_speaker_notes(
        slide,
        "~3 min | Allow 2-3 minutes for open discussion before advancing."
    )


def _add_learning_outcomes_slide(prs, slide_data, session_id, total_slides):
    """Learning Outcomes slide — achievements + areas for improvement."""
    slide = _new_slide(prs)
    slide_num = slide_data.get('slide_number', 0)
    content = slide_data.get('content', {})

    _add_title_text(slide, slide_data.get('title', 'Key Learning Points'))
    add_section_divider(slide, 1.6)

    achievements = content.get('key_achievements', [])
    areas = content.get('areas_for_improvement', [])

    # Two-column layout
    if achievements:
        items_text = "\n".join(f"\u2713 {a}" for a in achievements)
        col_width = 4.0 if areas else CONTENT_WIDTH
        add_feedback_box(
            slide, items_text, "KEY ACHIEVEMENTS", GREEN,
            Inches(MARGIN), Inches(2.0), Inches(col_width),
            Inches(min(1.2 + 0.35 * len(achievements), 3.5))
        )

    if areas:
        items_text = "\n".join(f"\u2192 {a}" for a in areas)
        left_x = 5.2 if achievements else MARGIN
        col_width = 4.0 if achievements else CONTENT_WIDTH
        add_feedback_box(
            slide, items_text, "AREAS FOR IMPROVEMENT", AMBER,
            Inches(left_x), Inches(2.0), Inches(col_width),
            Inches(min(1.2 + 0.35 * len(areas), 3.5))
        )

    # Debrief summary footer
    debrief_summary = content.get('debrief_summary', '')
    if debrief_summary:
        ds_box = slide.shapes.add_textbox(
            Inches(MARGIN), Inches(6.0), Inches(CONTENT_WIDTH), Inches(0.6)
        )
        dsp = ds_box.text_frame.paragraphs[0]
        dsp.text = debrief_summary
        dsp.font.size = Pt(14)
        dsp.font.italic = True
        dsp.font.color.rgb = SECONDARY
        dsp.alignment = PP_ALIGN.CENTER
        ds_box.text_frame.word_wrap = True

    add_slide_footer(slide, session_id, slide_num, total_slides)
    add_speaker_notes(
        slide,
        "~2 min | Summarize key points. Ask student to identify their top takeaway."
    )


def _add_action_plan_slide(prs, slide_data, session_id, total_slides):
    """Action Plan slide — numbered next steps."""
    slide = _new_slide(prs)
    slide_num = slide_data.get('slide_number', 0)
    content = slide_data.get('content', {})

    _add_title_text(slide, slide_data.get('title', 'Action Plan'))

    subtitle = slide_data.get('subtitle', '')
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(MARGIN), Inches(1.5), Inches(CONTENT_WIDTH), Inches(0.3)
        )
        sub_box.text_frame.paragraphs[0].text = _safe_text(subtitle)
        sub_box.text_frame.paragraphs[0].font.size = Pt(16)
        sub_box.text_frame.paragraphs[0].font.color.rgb = SECONDARY

    add_section_divider(slide, 1.85)

    # Action items
    actions = content.get('action_items', content.get('action_plan', []))
    y_pos = 2.2
    for idx, action in enumerate(actions[:6], 1):
        _add_numbered_badge(slide, idx, MARGIN, y_pos)

        action_box = slide.shapes.add_textbox(
            Inches(MARGIN + 0.55), Inches(y_pos), Inches(5.0), Inches(0.5)
        )
        atf = action_box.text_frame
        atf.word_wrap = True
        ap = atf.paragraphs[0]
        ap.text = _safe_text(action)
        ap.font.size = Pt(16)
        ap.font.color.rgb = BODY_TEXT

        y_pos += 0.6

    # Recommended focus areas
    focus = content.get('recommended_focus', [])
    if focus:
        focus_header = slide.shapes.add_textbox(
            Inches(MARGIN), Inches(y_pos + 0.3), Inches(CONTENT_WIDTH), Inches(0.3)
        )
        focus_header.text_frame.paragraphs[0].text = "RECOMMENDED FOCUS AREAS"
        focus_header.text_frame.paragraphs[0].font.size = Pt(10)
        focus_header.text_frame.paragraphs[0].font.bold = True
        focus_header.text_frame.paragraphs[0].font.color.rgb = SECONDARY

        focus_y = y_pos + 0.7
        for f_item in focus[:4]:
            f_box = slide.shapes.add_textbox(
                Inches(MARGIN + 0.2), Inches(focus_y), Inches(CONTENT_WIDTH - 0.4), Inches(0.35)
            )
            fp = f_box.text_frame.paragraphs[0]
            fp.text = f"\u2022 {_safe_text(f_item)}"
            fp.font.size = Pt(14)
            fp.font.color.rgb = PURPLE
            focus_y += 0.4

    add_slide_footer(slide, session_id, slide_num, total_slides)
    add_speaker_notes(
        slide,
        "~1 min | Agree on specific follow-up actions and timeline."
    )


def _add_discussion_slide(prs, slide_data, session_id, total_slides):
    """Discussion & Q&A slide — reflective prompts."""
    slide = _new_slide(prs)
    slide_num = slide_data.get('slide_number', 0)
    content = slide_data.get('content', {})

    # Centered title
    title_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(1.5), Inches(7), Inches(0.8)
    )
    tp = title_box.text_frame.paragraphs[0]
    tp.text = slide_data.get('title', 'Discussion')
    tp.font.size = Pt(40)
    tp.font.bold = True
    tp.font.color.rgb = NAVY
    tp.alignment = PP_ALIGN.CENTER

    subtitle = slide_data.get('subtitle', 'Questions & Reflection')
    sub_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.3), Inches(7), Inches(0.4)
    )
    sub_p = sub_box.text_frame.paragraphs[0]
    sub_p.text = subtitle
    sub_p.font.size = Pt(18)
    sub_p.font.color.rgb = SECONDARY
    sub_p.alignment = PP_ALIGN.CENTER

    # Prompts in 2x2 grid
    prompts = content.get('prompts', [
        "What was the most challenging aspect of this scenario?",
        "What would you do differently next time?",
        "How will you apply these learnings in practice?",
        "Any questions about the feedback provided?"
    ])

    positions = [
        (1.5, 3.3), (5.3, 3.3),
        (1.5, 4.8), (5.3, 4.8),
    ]

    for i, prompt in enumerate(prompts[:4]):
        px, py = positions[i]
        prompt_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(px), Inches(py), Inches(3.3), Inches(1.2)
        )
        prompt_shape.fill.solid()
        prompt_shape.fill.fore_color.rgb = CARD_FILL
        prompt_shape.line.color.rgb = PURPLE
        prompt_shape.line.width = Pt(2)

        ptf = prompt_shape.text_frame
        ptf.word_wrap = True
        ptf.margin_left = Inches(0.15)
        ptf.margin_right = Inches(0.15)
        ptf.margin_top = Inches(0.1)
        pp = ptf.paragraphs[0]
        pp.text = f'"{_safe_text(prompt)}"'
        pp.font.size = Pt(13)
        pp.font.italic = True
        pp.font.color.rgb = BODY_TEXT
        pp.alignment = PP_ALIGN.CENTER

    add_slide_footer(slide, session_id, slide_num, total_slides)
    add_speaker_notes(
        slide,
        "~5 min | Open floor for discussion. Allow the student to lead reflection."
    )


def _add_summary_slide(prs, slide_data, session_id, total_slides):
    """Summary & Close slide — metrics + key takeaways + references."""
    slide = _new_slide(prs)
    slide_num = slide_data.get('slide_number', 0)
    content = slide_data.get('content', {})

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(0.8), Inches(8), Inches(0.8)
    )
    tp = title_box.text_frame.paragraphs[0]
    tp.text = slide_data.get('title', 'Session Summary')
    tp.font.size = Pt(40)
    tp.font.bold = True
    tp.font.color.rgb = NAVY
    tp.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.5), Inches(8), Inches(0.4)
    )
    sp = sub_box.text_frame.paragraphs[0]
    sp.text = "Session Debrief Complete"
    sp.font.size = Pt(16)
    sp.font.color.rgb = SECONDARY
    sp.alignment = PP_ALIGN.CENTER

    add_section_divider(slide, 2.0)

    # Debrief summary text
    debrief_summary = content.get('debrief_summary', '')
    if debrief_summary:
        ds_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(2.3), Inches(7), Inches(0.6)
        )
        dsp = ds_box.text_frame.paragraphs[0]
        dsp.text = debrief_summary
        dsp.font.size = Pt(16)
        dsp.font.italic = True
        dsp.font.color.rgb = BODY_TEXT
        dsp.alignment = PP_ALIGN.CENTER
        ds_box.text_frame.word_wrap = True

    # Metric boxes row
    metrics_y = 3.2
    metric_data = []

    moments_reviewed = content.get('total_moments_reviewed')
    if moments_reviewed is not None:
        metric_data.append((str(moments_reviewed), "Moments Reviewed", CARD_FILL))

    moments_evaluated = content.get('total_moments_evaluated')
    if moments_evaluated is not None:
        metric_data.append((str(moments_evaluated), "Evaluated", RGBColor(220, 252, 231)))

    achievements = content.get('key_achievements')
    if isinstance(achievements, list):
        metric_data.append((str(len(achievements)), "Achievements", RGBColor(220, 252, 231)))
    elif isinstance(achievements, (int, float)):
        metric_data.append((str(achievements), "Achievements", RGBColor(220, 252, 231)))

    evidence_count = content.get('evidence_count')
    if evidence_count is not None:
        metric_data.append((str(evidence_count), "Evidence Links", RGBColor(219, 234, 254)))

    if metric_data:
        num_metrics = min(len(metric_data), 4)
        m_width = (CONTENT_WIDTH - 0.3 * (num_metrics - 1)) / num_metrics
        for i, (val, label, bg) in enumerate(metric_data[:4]):
            mx = MARGIN + i * (m_width + 0.3)
            _add_metric_box(slide, val, label, bg, mx, metrics_y, width=m_width)

    # Strength/Growth Ratio Bar
    strength_count = content.get('strength_count', 0)
    growth_count = content.get('growth_count', 0)
    total_feedback = strength_count + growth_count
    if total_feedback > 0:
        bar_y = 4.5
        bar_width = 6.0
        bar_height = 0.3
        bar_left = MARGIN + (CONTENT_WIDTH - bar_width) / 2

        # Label
        ratio_label = slide.shapes.add_textbox(
            Inches(bar_left), Inches(bar_y - 0.3), Inches(bar_width), Inches(0.25)
        )
        rlp = ratio_label.text_frame.paragraphs[0]
        rlp.text = "FEEDBACK BALANCE"
        rlp.font.size = Pt(10)
        rlp.font.bold = True
        rlp.font.color.rgb = SECONDARY
        rlp.alignment = PP_ALIGN.CENTER

        # Green bar (strengths)
        green_width = (strength_count / total_feedback) * bar_width
        if green_width > 0:
            green_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(bar_left), Inches(bar_y),
                Inches(green_width), Inches(bar_height)
            )
            green_bar.fill.solid()
            green_bar.fill.fore_color.rgb = GREEN
            green_bar.line.fill.background()

        # Amber bar (growth)
        amber_width = bar_width - green_width
        if amber_width > 0:
            amber_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(bar_left + green_width), Inches(bar_y),
                Inches(amber_width), Inches(bar_height)
            )
            amber_bar.fill.solid()
            amber_bar.fill.fore_color.rgb = AMBER
            amber_bar.line.fill.background()

        # Labels below bar
        count_box = slide.shapes.add_textbox(
            Inches(bar_left), Inches(bar_y + 0.32), Inches(bar_width), Inches(0.25)
        )
        ctf = count_box.text_frame
        cp = ctf.paragraphs[0]
        sr = cp.add_run()
        sr.text = f"\u2713 {strength_count} Strengths"
        sr.font.size = Pt(10)
        sr.font.color.rgb = GREEN
        sr.font.bold = True
        sp_run = cp.add_run()
        sp_run.text = "    "
        sp_run.font.size = Pt(10)
        gr = cp.add_run()
        gr.text = f"\u2192 {growth_count} Growth Areas"
        gr.font.size = Pt(10)
        gr.font.color.rgb = AMBER
        gr.font.bold = True
        cp.alignment = PP_ALIGN.CENTER

    # Key takeaways
    takeaways = slide_data.get('key_takeaways', [])
    tk_start_y = 5.3 if total_feedback > 0 else 4.7
    if takeaways:
        tk_header = slide.shapes.add_textbox(
            Inches(MARGIN), Inches(tk_start_y), Inches(CONTENT_WIDTH), Inches(0.3)
        )
        tk_header.text_frame.paragraphs[0].text = "KEY TAKEAWAYS"
        tk_header.text_frame.paragraphs[0].font.size = Pt(10)
        tk_header.text_frame.paragraphs[0].font.bold = True
        tk_header.text_frame.paragraphs[0].font.color.rgb = SECONDARY

        tk_y = tk_start_y + 0.4
        for idx, takeaway in enumerate(takeaways[:4], 1):
            _add_numbered_badge(slide, idx, MARGIN, tk_y)
            tk_box = slide.shapes.add_textbox(
                Inches(MARGIN + 0.5), Inches(tk_y), Inches(CONTENT_WIDTH - 0.6), Inches(0.5)
            )
            tkp = tk_box.text_frame.paragraphs[0]
            tkp.text = _safe_text(takeaway)
            tkp.font.size = Pt(14)
            tkp.font.color.rgb = BODY_TEXT
            tk_box.text_frame.word_wrap = True
            tk_y += 0.55

    # References
    references = content.get('references', [])
    if references:
        ref_y = max(tk_y + 0.3 if takeaways else 4.7, 6.0)
        ref_box = slide.shapes.add_textbox(
            Inches(MARGIN), Inches(ref_y), Inches(CONTENT_WIDTH), Inches(0.8)
        )
        rtf = ref_box.text_frame
        rtf.word_wrap = True
        rp = rtf.paragraphs[0]
        rp.text = "REFERENCES"
        rp.font.size = Pt(9)
        rp.font.bold = True
        rp.font.color.rgb = SECONDARY

        for ref in references[:3]:
            citation = ref.get('citation', '') if isinstance(ref, dict) else str(ref)
            rp2 = rtf.add_paragraph()
            rp2.text = f"\u2022 {_safe_text(citation)}"
            rp2.font.size = Pt(10)
            rp2.font.color.rgb = SECONDARY

    add_slide_footer(slide, session_id, slide_num, total_slides)
    add_speaker_notes(
        slide,
        "Closing remarks. Provide a copy of this presentation to the student."
    )
